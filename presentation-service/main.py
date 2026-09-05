import os
import time
import json
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional

from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Header, Depends, BackgroundTasks, Request, status
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

if os.getenv("APP_ENV", "development") != "production":
    load_dotenv(Path(__file__).resolve().parent.parent / ".env")

def log_event(event: str, **fields: object) -> None:
    print(json.dumps({"service": "presentation-service", "event": event, **fields}), flush=True)


INTERNAL_SERVICE_KEY = os.getenv("INTERNAL_SERVICE_KEY", "")
if not INTERNAL_SERVICE_KEY:
    raise RuntimeError("INTERNAL_SERVICE_KEY environment variable is required.")

app = FastAPI(title="Web Analytics - Presentation Engine", version="1.0.0")

@app.middleware("http")
async def log_requests(request: Request, call_next):
    request_id = request.headers.get("X-Request-Id") or str(uuid.uuid4())
    started_at = time.perf_counter()
    response = await call_next(request)
    duration_ms = round((time.perf_counter() - started_at) * 1000)
    response.headers["X-Request-Id"] = request_id
    log_event("http.request.completed", requestId=request_id, method=request.method, path=request.url.path, status=response.status_code, durationMs=duration_ms)
    return response


TEMP_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "temp_outputs")
os.makedirs(TEMP_DIR, exist_ok=True)

def verify_internal_key(x_internal_key: str = Header(None, alias="X-Internal-Key")):
    if not x_internal_key or x_internal_key != INTERNAL_SERVICE_KEY:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid or missing internal service authorization key"
        )

def cleanup_file(path: str):
    try:
        if os.path.exists(path):
            os.remove(path)
    except Exception as error:
        print(f"[Cleanup Warning] Failed to delete {path}: {error}")

def sweep_temp_dir():
    now = time.time()
    one_hour_ago = now - 3600
    try:
        for fname in os.listdir(TEMP_DIR):
            fpath = os.path.join(TEMP_DIR, fname)
            if os.path.isfile(fpath) and os.path.getmtime(fpath) < one_hour_ago:
                try:
                    os.remove(fpath)
                except Exception:
                    pass
    except Exception as error:
        print(f"[Sweep Warning] Error cleaning temp directory: {error}")

@app.on_event("startup")
def on_startup():
    sweep_temp_dir()

class DeckRequest(BaseModel):
    topic: str = Field(default="Market Research")
    bullets: List[str] = Field(default_factory=list)
    metrics: Dict[str, int] = Field(default_factory=lambda: {"Positive": 15, "Negative": 7, "Neutral": 11})
    sources: Optional[List[Dict[str, Any]]] = Field(default_factory=list)

def generate_analytics_chart(data_summary: Dict[str, int], output_img_path: str):
    if not data_summary:
        data_summary = {"Positive": 10, "Negative": 5, "Neutral": 8}

    labels = list(data_summary.keys())
    values = list(data_summary.values())

    palette = ['#1E3A8A', '#0D9488', '#F59E0B', '#EF4444', '#6366F1']
    colors = [palette[i % len(palette)] for i in range(len(labels))]

    fig, ax = plt.subplots(figsize=(7, 4.2), dpi=220)
    fig.patch.set_facecolor('#F8FAFC')
    ax.set_facecolor('#F8FAFC')

    bars = ax.bar(labels, values, color=colors, width=0.55, edgecolor="#0F172A", linewidth=0.6)

    for bar in bars:
        height = bar.get_height()
        ax.annotate(
            f'{height}',
            xy=(bar.get_x() + bar.get_width() / 2, height),
            xytext=(0, 4),
            textcoords="offset points",
            ha='center',
            va='bottom',
            fontsize=10,
            fontweight='bold',
            color='#1E293B'
        )

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CBD5E1')
    ax.spines['bottom'].set_color('#94A3B8')

    plt.title("Search Result Sentiment Summary", fontsize=13, fontweight='bold', color='#0F172A', pad=15)
    plt.ylabel("Result Count", fontsize=10, color='#475569', labelpad=8)
    plt.xticks(rotation=15, ha='right', fontsize=9, color='#1E293B')
    plt.yticks(color='#475569')
    plt.grid(axis='y', linestyle='--', alpha=0.3, color='#94A3B8')
    plt.tight_layout()
    plt.savefig(output_img_path, dpi=220, facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close()

@app.get("/health")
def health_check():
    return {"status": "ok", "service": "presentation-service", "version": "1.0.0"}

@app.post("/generate-presentation", dependencies=[Depends(verify_internal_key)])
def make_deck(payload: DeckRequest, background_tasks: BackgroundTasks):
    topic = payload.topic
    bullet_points = payload.bullets
    chart_metrics = payload.metrics
    sources = payload.sources or []

    run_id = str(uuid.uuid4())[:8]
    clean_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()
    filename_base = f"Research_Report_{clean_topic.replace(' ', '_')}_{run_id}"
    ppt_path = os.path.join(TEMP_DIR, f"{filename_base}.pptx")
    chart_path = os.path.join(TEMP_DIR, f"temp_chart_{run_id}.png")

    try:
        generate_analytics_chart(chart_metrics, chart_path)

        prs = Presentation()

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Research Summary"
        slide.placeholders[1].text = f"Web Analysis Report\nTopic: {topic}"

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Key Findings"

        text_frame = slide.placeholders[1].text_frame
        text_frame.word_wrap = True

        effective_bullets = bullet_points if bullet_points else [
            f"Retrieved live search-result snippets for '{topic}'.",
            "Classified title and snippet text with a lexical sentiment heuristic.",
            "Counted keyword hits in titles and snippets for the requested terms.",
            "Included the retrieved source URL and snippet for each result card."
        ]

        for index, bullet in enumerate(effective_bullets[:5]):
            paragraph = text_frame.add_paragraph() if index > 0 else text_frame.paragraphs[0]
            paragraph.text = f"- {bullet}"
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = RGBColor(30, 41, 59)

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        text_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.8))
        chart_text_frame = text_box.text_frame
        title_paragraph = chart_text_frame.paragraphs[0]
        title_paragraph.text = "Sentiment counts"
        title_paragraph.font.size = Pt(22)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

        subtitle = chart_text_frame.add_paragraph()
        subtitle.text = f"Counts derived from lexical heuristics over retrieved titles and snippets for: {topic}"
        subtitle.font.size = Pt(12)
        subtitle.font.color.rgb = RGBColor(100, 116, 139)

        if os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(1.8), Inches(1.6), width=Inches(6.4))

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sources"

        rows = min(len(sources) + 1, 6)
        cols = 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.8), Inches(3.2))
        table = table_shape.table

        table.columns[0].width = Inches(2.8)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(1.5)

        headers = ["Source / Title", "Snippet", "Sentiment"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 58, 138)

        for row_index in range(1, rows):
            source = sources[row_index - 1] if row_index - 1 < len(sources) else {}
            title_val = (source.get("title") or source.get("sourceUrl") or f"Source {row_index}")[:40]
            snippet_text = source.get("snippet") or "Retrieved result snippet unavailable."
            snippet_val = snippet_text[:80] + ("..." if len(snippet_text) > 80 else "")
            sentiment_val = (source.get("sentiment") or "NEUTRAL").upper()

            row_data = [title_val, snippet_val, sentiment_val]
            for col_index, value in enumerate(row_data):
                cell = table.cell(row_index, col_index)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = RGBColor(51, 65, 85)

        prs.save(ppt_path)

        if os.path.exists(chart_path):
            try:
                os.remove(chart_path)
            except Exception:
                pass

        background_tasks.add_task(cleanup_file, ppt_path)

        return FileResponse(
            ppt_path,
            filename=f"{filename_base}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as error:
        print(f"[Presentation Error] Failed to generate deck: {error}")
        if os.path.exists(chart_path):
            try:
                os.remove(chart_path)
            except Exception:
                pass
        if os.path.exists(ppt_path):
            try:
                os.remove(ppt_path)
            except Exception:
                pass
        raise HTTPException(status_code=500, detail="Presentation generation failed")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="127.0.0.1", port=8002, reload=True)
